"""
文档加载器 - 支持多种文档格式
支持：PDF、Word、Excel、PPT、Markdown、TXT、代码文件、电子书等
"""
import os
import re
import csv
import io
from pathlib import Path
from typing import List, Dict, Any, Optional, Union
from dataclasses import dataclass


@dataclass
class Document:
    """文档对象"""
    content: str
    metadata: Dict[str, Any]
    source: str  # 文件路径或来源标识


class DocumentLoader:
    """文档加载器基类"""

    SUPPORTED_EXTENSIONS = {
        # 办公文档
        '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.csv',
        '.pptx', '.ppt', '.odt', '.ods', '.odp',
        # 文本与 Markdown
        '.txt', '.md', '.markdown', '.rtf',
        # 代码文件
        '.py', '.js', '.ts', '.jsx', '.tsx',
        '.java', '.cpp', '.c', '.h', '.hpp',
        '.go', '.rs', '.rb', '.php',
        '.html', '.css', '.json', '.xml', '.yaml', '.yml',
        '.sql', '.sh', '.bat', '.ps1', '.swift', '.kt',
        # 电子书
        '.epub',
    }

    def __init__(self, encoding: str = 'utf-8'):
        self.encoding = encoding

    def load(self, source: Union[str, Path]) -> List[Document]:
        """加载文档"""
        source = Path(source)

        if source.is_file():
            return self._load_file(source)
        elif source.is_dir():
            return self._load_directory(source)
        else:
            raise ValueError(f"路径不存在: {source}")

    def _load_file(self, file_path: Path) -> List[Document]:
        """加载单个文件"""
        suffix = file_path.suffix.lower()

        if suffix == '.pdf':
            return self._load_pdf(file_path)
        elif suffix in ['.docx']:
            return self._load_docx(file_path)
        elif suffix in ['.doc']:
            return self._load_doc(file_path)
        elif suffix in ['.xlsx']:
            return self._load_xlsx(file_path)
        elif suffix in ['.xls']:
            return self._load_xls(file_path)
        elif suffix == '.csv':
            return self._load_csv(file_path)
        elif suffix in ['.pptx']:
            return self._load_pptx(file_path)
        elif suffix in ['.ppt']:
            return self._load_ppt(file_path)
        elif suffix in ['.odt']:
            return self._load_odt(file_path)
        elif suffix in ['.ods']:
            return self._load_ods(file_path)
        elif suffix in ['.odp']:
            return self._load_odp(file_path)
        elif suffix == '.epub':
            return self._load_epub(file_path)
        elif suffix in ['.md', '.markdown']:
            return self._load_markdown(file_path)
        elif suffix in ['.rtf']:
            return self._load_rtf(file_path)
        elif suffix in self.SUPPORTED_EXTENSIONS:
            return self._load_text(file_path)
        else:
            # 尝试作为文本文件加载
            try:
                return self._load_text(file_path)
            except Exception as e:
                print(f"⚠️ 无法加载文件 {file_path}: {e}")
                return []

    def _load_directory(self, dir_path: Path) -> List[Document]:
        """加载目录下所有支持的文件"""
        documents = []

        for file_path in dir_path.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS:
                try:
                    docs = self._load_file(file_path)
                    documents.extend(docs)
                    print(f"✅ 已加载: {file_path.relative_to(dir_path)}")
                except Exception as e:
                    print(f"⚠️ 加载失败 {file_path}: {e}")

        return documents

    def _load_text(self, file_path: Path) -> List[Document]:
        """加载文本文件"""
        try:
            with open(file_path, 'r', encoding=self.encoding) as f:
                content = f.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            with open(file_path, 'r', encoding='gbk') as f:
                content = f.read()

        # 清理内容
        content = self._clean_text(content)

        metadata = {
            'source': str(file_path),
            'filename': file_path.name,
            'file_type': file_path.suffix,
            'file_size': len(content),
        }

        return [Document(content=content, metadata=metadata, source=str(file_path))]

    def _load_markdown(self, file_path: Path) -> List[Document]:
        """加载 Markdown 文件，提取结构化信息"""
        docs = self._load_text(file_path)

        if docs:
            # 提取标题结构
            content = docs[0].content
            headers = self._extract_markdown_headers(content)
            docs[0].metadata['headers'] = headers
            docs[0].metadata['file_type'] = 'markdown'

        return docs

    def _load_pdf(self, file_path: Path) -> List[Document]:
        """加载 PDF 文件"""
        try:
            import PyPDF2
        except ImportError:
            print("⚠️ PyPDF2 未安装，尝试使用 pip install PyPDF2")
            return []

        documents = []

        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            num_pages = len(pdf_reader.pages)

            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    metadata = {
                        'source': str(file_path),
                        'filename': file_path.name,
                        'file_type': 'pdf',
                        'page': page_num,
                        'total_pages': num_pages,
                    }
                    documents.append(Document(
                        content=text,
                        metadata=metadata,
                        source=str(file_path)
                    ))

        return documents

    def _load_docx(self, file_path: Path) -> List[Document]:
        """加载 Word (.docx) 文件"""
        try:
            import docx
        except ImportError:
            print("⚠️ python-docx 未安装，尝试使用 pip install python-docx")
            return []

        doc = docx.Document(file_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        content = '\n'.join(paragraphs)
        content = self._clean_text(content)

        metadata = {
            'source': str(file_path),
            'filename': file_path.name,
            'file_type': 'docx',
            'paragraph_count': len(paragraphs),
        }
        return [Document(content=content, metadata=metadata, source=str(file_path))]

    def _load_doc(self, file_path: Path) -> List[Document]:
        """加载旧版 Word (.doc) 文件"""
        # 尝试使用 antiword 或 catdoc 等外部工具
        try:
            import subprocess
            result = subprocess.run(
                ['antiword', str(file_path)],
                capture_output=True,
                text=True,
                check=True
            )
            content = self._clean_text(result.stdout)
        except Exception:
            try:
                result = subprocess.run(
                    ['catdoc', str(file_path)],
                    capture_output=True,
                    text=True,
                    check=True
                )
                content = self._clean_text(result.stdout)
            except Exception:
                print(f"⚠️ 无法加载 .doc 文件 {file_path}，请安装 antiword 或 catdoc")
                return []

        metadata = {
            'source': str(file_path),
            'filename': file_path.name,
            'file_type': 'doc',
        }
        return [Document(content=content, metadata=metadata, source=str(file_path))]

    def _load_xlsx(self, file_path: Path) -> List[Document]:
        """加载 Excel (.xlsx) 文件"""
        try:
            import openpyxl
        except ImportError:
            print("⚠️ openpyxl 未安装，尝试使用 pip install openpyxl")
            return []

        wb = openpyxl.load_workbook(file_path, data_only=True)
        sheets_content = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                sheets_content.append(f"【Sheet: {sheet_name}】\n" + '\n'.join(rows))

        content = '\n\n'.join(sheets_content)
        content = self._clean_text(content)

        metadata = {
            'source': str(file_path),
            'filename': file_path.name,
            'file_type': 'xlsx',
            'sheet_count': len(wb.sheetnames),
        }
        return [Document(content=content, metadata=metadata, source=str(file_path))]

    def _load_xls(self, file_path: Path) -> List[Document]:
        """加载旧版 Excel (.xls) 文件"""
        try:
            import xlrd
        except ImportError:
            print("⚠️ xlrd 未安装，尝试使用 pip install xlrd")
            return []

        wb = xlrd.open_workbook(file_path)
        sheets_content = []
        for sheet_idx in range(wb.nsheets):
            sheet = wb.sheet_by_index(sheet_idx)
            rows = []
            for row_idx in range(sheet.nrows):
                row_values = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
                row_text = ' | '.join(row_values)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                sheets_content.append(f"【Sheet: {sheet.name}】\n" + '\n'.join(rows))

        content = '\n\n'.join(sheets_content)
        content = self._clean_text(content)

        metadata = {
            'source': str(file_path),
            'filename': file_path.name,
            'file_type': 'xls',
            'sheet_count': wb.nsheets,
        }
        return [Document(content=content, metadata=metadata, source=str(file_path))]

    def _load_csv(self, file_path: Path) -> List[Document]:
        """加载 CSV 文件"""
        try:
            with open(file_path, 'r', encoding=self.encoding, newline='') as f:
                reader = csv.reader(f)
                rows = [' | '.join(row) for row in reader]
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='gbk', newline='') as f:
                reader = csv.reader(f)
                rows = [' | '.join(row) for row in reader]

        content = '\n'.join(rows)
        content = self._clean_text(content)

        metadata = {
            'source': str(file_path),
            'filename': file_path.name,
            'file_type': 'csv',
            'row_count': len(rows),
        }
        return [Document(content=content, metadata=metadata, source=str(file_path))]

    def _load_pptx(self, file_path: Path) -> List[Document]:
        """加载 PowerPoint (.pptx) 文件"""
        try:
            from pptx import Presentation
        except ImportError:
            print("⚠️ python-pptx 未安装，尝试使用 pip install python-pptx")
            return []

        prs = Presentation(file_path)
        slides_content = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides_content.append(f"【第 {i} 页】\n" + '\n'.join(texts))

        content = '\n\n'.join(slides_content)
        content = self._clean_text(content)

        metadata = {
            'source': str(file_path),
            'filename': file_path.name,
            'file_type': 'pptx',
            'slide_count': len(prs.slides),
        }
        return [Document(content=content, metadata=metadata, source=str(file_path))]

    def _load_ppt(self, file_path: Path) -> List[Document]:
        """加载旧版 PowerPoint (.ppt) 文件"""
        try:
            import subprocess
            result = subprocess.run(
                ['catppt', str(file_path)],
                capture_output=True,
                text=True,
                check=True
            )
            content = self._clean_text(result.stdout)
        except Exception:
            print(f"⚠️ 无法加载 .ppt 文件 {file_path}，请安装 catppt")
            return []

        metadata = {
            'source': str(file_path),
            'filename': file_path.name,
            'file_type': 'ppt',
        }
        return [Document(content=content, metadata=metadata, source=str(file_path))]

    def _load_odt(self, file_path: Path) -> List[Document]:
        """加载 OpenDocument Text (.odt) 文件"""
        try:
            from odf import opendocument
            from odf.text import P
            doc = opendocument.load(file_path)
            paragraphs = [str(p) for p in doc.spreadsheet.getElementsByType(P) if str(p).strip()]
            content = '\n'.join(paragraphs)
            content = self._clean_text(content)
        except Exception as e:
            print(f"⚠️ 无法加载 .odt 文件 {file_path}: {e}")
            return []

        metadata = {
            'source': str(file_path),
            'filename': file_path.name,
            'file_type': 'odt',
        }
        return [Document(content=content, metadata=metadata, source=str(file_path))]

    def _load_ods(self, file_path: Path) -> List[Document]:
        """加载 OpenDocument Spreadsheet (.ods) 文件"""
        try:
            import pandas as pd
            dfs = pd.read_excel(file_path, sheet_name=None)
            sheets_content = []
            for sheet_name, df in dfs.items():
                rows = [' | '.join(str(cell) for cell in row) for row in df.values]
                if rows:
                    sheets_content.append(f"【Sheet: {sheet_name}】\n" + '\n'.join(rows))
            content = '\n\n'.join(sheets_content)
            content = self._clean_text(content)
        except Exception as e:
            print(f"⚠️ 无法加载 .ods 文件 {file_path}: {e}")
            return []

        metadata = {
            'source': str(file_path),
            'filename': file_path.name,
            'file_type': 'ods',
        }
        return [Document(content=content, metadata=metadata, source=str(file_path))]

    def _load_odp(self, file_path: Path) -> List[Document]:
        """加载 OpenDocument Presentation (.odp) 文件"""
        try:
            from odf import opendocument
            from odf.draw import Page
            doc = opendocument.load(file_path)
            pages = doc.presentation.getElementsByType(Page)
            pages_content = []
            for i, page in enumerate(pages, 1):
                texts = []
                for node in page.childNodes:
                    text = str(node)
                    if text.strip():
                        texts.append(text.strip())
                if texts:
                    pages_content.append(f"【第 {i} 页】\n" + '\n'.join(texts))
            content = '\n\n'.join(pages_content)
            content = self._clean_text(content)
        except Exception as e:
            print(f"⚠️ 无法加载 .odp 文件 {file_path}: {e}")
            return []

        metadata = {
            'source': str(file_path),
            'filename': file_path.name,
            'file_type': 'odp',
        }
        return [Document(content=content, metadata=metadata, source=str(file_path))]

    def _load_epub(self, file_path: Path) -> List[Document]:
        """加载 EPUB 电子书"""
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
        except ImportError:
            print("⚠️ ebooklib 或 beautifulsoup4 未安装，尝试使用 pip install ebooklib beautifulsoup4")
            return []

        book = epub.read_epub(file_path)
        chapters = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text = soup.get_text(separator='\n').strip()
                if text:
                    chapters.append(text)

        content = '\n\n'.join(chapters)
        content = self._clean_text(content)

        metadata = {
            'source': str(file_path),
            'filename': file_path.name,
            'file_type': 'epub',
            'chapter_count': len(chapters),
        }
        return [Document(content=content, metadata=metadata, source=str(file_path))]

    def _load_rtf(self, file_path: Path) -> List[Document]:
        """加载 RTF 文件"""
        try:
            import striprtf
            from striprtf.striprtf import rtf_to_text
        except ImportError:
            # 降级方案：尝试作为文本读取并清理 RTF 标记
            try:
                with open(file_path, 'r', encoding=self.encoding) as f:
                    content = f.read()
                # 简单去除 RTF 控制字
                content = re.sub(r'\\[a-z]+\d*\s?', ' ', content)
                content = re.sub(r'[{}]', '', content)
                content = self._clean_text(content)
            except Exception as e:
                print(f"⚠️ 无法加载 RTF 文件 {file_path}: {e}")
                return []
        else:
            with open(file_path, 'r', encoding=self.encoding) as f:
                content = rtf_to_text(f.read())
            content = self._clean_text(content)

        metadata = {
            'source': str(file_path),
            'filename': file_path.name,
            'file_type': 'rtf',
        }
        return [Document(content=content, metadata=metadata, source=str(file_path))]

    def _clean_text(self, text: str) -> str:
        """清理文本内容"""
        # 移除多余的空白字符
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = re.sub(r'[ \t]+', ' ', text)
        # 移除特殊控制字符
        text = re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f]', '', text)
        return text.strip()

    def _extract_markdown_headers(self, content: str) -> List[Dict[str, Any]]:
        """提取 Markdown 标题结构"""
        headers = []
        for match in re.finditer(r'^(#{1,6})\s+(.+)$', content, re.MULTILINE):
            level = len(match.group(1))
            title = match.group(2).strip()
            headers.append({'level': level, 'title': title})
        return headers


def load_documents(
    source: Union[str, Path],
    encoding: str = 'utf-8',
    file_filter: Optional[List[str]] = None
) -> List[Document]:
    """
    便捷函数：加载文档

    Args:
        source: 文件或目录路径
        encoding: 文件编码
        file_filter: 文件扩展名过滤器，如 ['.md', '.txt']

    Returns:
        Document 列表
    """
    loader = DocumentLoader(encoding=encoding)

    # 如果有过滤器，临时修改支持的后缀列表
    if file_filter:
        loader.SUPPORTED_EXTENSIONS = set(file_filter)

    return loader.load(source)


# 便捷函数
def load_pdf(file_path: Union[str, Path]) -> List[Document]:
    """加载 PDF 文件"""
    loader = DocumentLoader()
    return loader._load_pdf(Path(file_path))


def load_markdown(file_path: Union[str, Path]) -> List[Document]:
    """加载 Markdown 文件"""
    loader = DocumentLoader()
    return loader._load_markdown(Path(file_path))


def load_text(file_path: Union[str, Path]) -> List[Document]:
    """加载文本文件"""
    loader = DocumentLoader()
    return loader._load_text(Path(file_path))
